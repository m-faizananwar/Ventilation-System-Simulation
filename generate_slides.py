import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --- CONFIGURATION ---
DARK_BG = RGBColor(10, 10, 15)
TEXT_MAIN = RGBColor(255, 255, 255)
TEXT_SEC = RGBColor(160, 160, 176)
ACCENT_GREEN = RGBColor(0, 255, 136)
ACCENT_BLUE = RGBColor(68, 136, 255)
ACCENT_RED = RGBColor(255, 68, 102)
ACCENT_PURPLE = RGBColor(153, 102, 255)
ACCENT_YELLOW = RGBColor(255, 204, 0)

def create_presentation():
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- HELPER FUNCTIONS ---
    def add_slide(layout_index=6): # 6 is blank
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        return slide

    def add_title(slide, text, font_size=40, top=Inches(0.5)):
        title_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_MAIN
        p.font.name = "Arial"
        p.font.bold = True
        return title_box

    def add_text_box(slide, text, left, top, width, height, font_size=18, color=TEXT_SEC, bold=False):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.font.bold = bold
        return tf

    def add_bullet(tf, text, level=0):
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SEC
        p.font.name = "Arial"

    def add_code_block(slide, code_text, left, top, width, height):
        # Background box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(26, 26, 46) # Darker code bg
        shape.line.color.rgb = RGBColor(50, 50, 60)
        
        # Text
        tf = shape.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = code_text
        p.font.name = "Courier New"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.LEFT

    # ================= SLIDE 1: TITLE =================
    s1 = add_slide()
    # Tag
    add_text_box(s1, "CS-221 PROJECT PRESENTATION", Inches(1), Inches(2), Inches(5), Inches(0.5), font_size=14, color=ACCENT_GREEN, bold=True)
    
    # Title
    t1 = add_text_box(s1, "Smart Corridor Ventilation &\nAir Quality Monitoring System", Inches(1), Inches(2.5), Inches(11), Inches(2), font_size=44, color=TEXT_MAIN, bold=True)
    
    # Subtitle
    add_text_box(s1, "A RISC-V based embedded system integrating IoT sensors, low-level architecture,\nand 3D visualization for autonomous indoor air quality control.", Inches(1), Inches(4.5), Inches(10), Inches(1), font_size=18, color=TEXT_SEC)
    
    # Info
    info_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.8), Inches(11.3), Inches(0.8))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RGBColor(25, 25, 35)
    info_box.line.color.rgb = RGBColor(60, 60, 70)
    
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CS-221  |  Computer Organization & Assembly Language  |  BSCS-3B  |  Group 05"
    p.font.color.rgb = TEXT_MAIN
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    info_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ================= SLIDE 2: TEAM =================
    s2 = add_slide()
    add_title(s2, "Meets the Team")
    
    members = [
        ("Faizan Anwar", "455259", "Team Rep • System Architect & Lead Developer", ACCENT_BLUE),
        ("Abdul Moiz", "465932", "Hardware Simulation Specialist", ACCENT_RED),
        ("Muhammad Taha", "467244", "Software Engineer (Assembly)", ACCENT_GREEN),
        ("Sham", "457919", "Testing & Validation Engineer", ACCENT_PURPLE),
    ]
    
    for i, (name, cms, role, color) in enumerate(members):
        left = Inches(1 + (i * 3))
        # Card bg
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(2.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(20, 20, 30)
        card.line.color.rgb = color
        
        # Initials circle
        circle = s2.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.9), Inches(3), Inches(1), Inches(1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        initials = "".join([n[0] for n in name.split()[:2]])
        p.text = initials
        p.font.bold = True
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        
        # Name
        tf_name = add_text_box(s2, name, left, Inches(4.2), Inches(2.8), Inches(0.5), font_size=16, color=TEXT_MAIN, bold=True)
        tf_name.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # CMS
        tf_cms = add_text_box(s2, cms, left, Inches(4.6), Inches(2.8), Inches(0.4), font_size=14, color=color, bold=False)
        tf_cms.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Role
        tf_role = add_text_box(s2, role, left + Inches(0.2), Inches(5), Inches(2.4), Inches(1), font_size=12, color=TEXT_SEC, bold=False)
        tf_role.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ================= SLIDE 3: PROBLEM STATEMENT =================
    s3 = add_slide()
    add_title(s3, "Problem Statement")
    
    # Box
    box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11), Inches(4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(20, 20, 30)
    box.line.color.rgb = ACCENT_RED
    
    tf = box.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.5)
    
    p = tf.paragraphs[0]
    p.text = "The Challenge"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "\nConventional ventilation systems are manual or timer-based, failing to respond dynamically to real-time hazards choices."
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MAIN
    
    bullets = [
        "Hypercapnia Risk: Elevated CO2 levels (>1000 ppm) cause fatigue",
        "Particulate Accumulation: PM2.5 particles (>100 µg/m³) damage health",
        "Fire Hazards: Smoke infiltration requires immediate mitigation",
        "Delayed Response: Manual systems have dangerous reaction latency"
    ]
    
    for b in bullets:
        add_bullet(tf, b)

    # ================= SLIDE 4: OBJECTIVES =================
    s4 = add_slide()
    add_title(s4, "Project Objectives")
    
    objs = [
        ("Real-Time Response", "Detect hazards within milliseconds using RISC-V", "⚡"),
        ("Autonomous Control", "Activate ventilation without human intervention", "🤖"),
        ("Multi-Sensor Fusion", "Combine Smoke, PM2.5, CO2, and Temperature", "📊"),
        ("Digital Twin", "Immersive 3D simulation for monitoring", "🏠")
    ]
    
    for i, (title, desc, icon) in enumerate(objs):
        x = Inches(1 + (i%2)*6)
        y = Inches(2 + (i//2)*2.5)
        
        # Box
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(20, 20, 30)
        box.line.color.rgb = ACCENT_BLUE

        # Icon (Text)
        add_text_box(s4, icon, x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(1), font_size=30)
        
        # Content
        add_text_box(s4, title, x + Inches(1.2), y + Inches(0.3), Inches(4), Inches(0.5), font_size=18, color=TEXT_MAIN, bold=True)
        add_text_box(s4, desc, x + Inches(1.2), y + Inches(0.8), Inches(4), Inches(1), font_size=14, color=TEXT_SEC)

    # ================= SLIDE 5: ARCHITECTURE OVERVIEW =================
    s5 = add_slide()
    add_title(s5, "System Architecture: Three-Layer Design")
    
    layers = [
        ("Layer 1: Ripes", "RISC-V Architecture", "Low-level processor simulation with Memory-Mapped I/O and register manipulation.", ACCENT_RED),
        ("Layer 2: Wokwi", "IoT Hardware Layer", "ESP32-based sensor node with real sensors, WiFi, and MQTT communication.", ACCENT_BLUE),
        ("Layer 3: Simulation", "3D Digital Twin", "React + Three.js immersive environment with particles and physics.", ACCENT_GREEN)
    ]
    
    for i, (l_name, l_sub, l_desc, color) in enumerate(layers):
        x = Inches(0.5 + i*4.2)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(4), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(20, 20, 30)
        box.line.color.rgb = color
        
        tf = box.text_frame
        tf.paragraphs[0].text = l_name
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = l_sub
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + l_desc
        p2.font.size = Pt(14)
        p2.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 6: BLOCK DIAGRAM =================
    s6 = add_slide()
    add_title(s6, "Architecture Block Diagram")
    
    # Draw simple blocks
    # CPU
    cpu = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(3.33), Inches(1))
    cpu.text_frame.text = "RISC-V Core\nRegisters x0-x31 | ALU"
    cpu.fill.solid()
    cpu.fill.fore_color.rgb = RGBColor(40, 20, 20)
    cpu.line.color.rgb = ACCENT_RED

    # Bus
    bus = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(11.33), Inches(0.5))
    bus.text_frame.text = "System Bus (Data / Address / Control)"
    bus.fill.solid()
    bus.fill.fore_color.rgb = RGBColor(20, 30, 50)
    bus.line.color.rgb = ACCENT_BLUE
    
    # Connection
    conn = s6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.66), Inches(2.5), Inches(6.66), Inches(3)) # simple representation

    # Memory Banks
    banks = [
        ("0xF0000000\nSmoke + PM2.5", Inches(1.5)),
        ("0xF0000004\nCO2 Level", Inches(5)),
        ("0xF0000008\nLED Matrix", Inches(8.5))
    ]
    
    for txt, x in banks:
        b = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4), Inches(2.5), Inches(1))
        b.text_frame.text = txt
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(20, 40, 30)
        b.line.color.rgb = ACCENT_GREEN

    # Peripherals
    peripherals = [
        ("Smoke Detector", Inches(1.5)),
        ("PM2.5 Sensor", Inches(3.5)),
        ("CO2 Sensor", Inches(5.5)),
        ("Fan + LEDs", Inches(8.5))
    ]
    
    for txt, x in peripherals:
        p = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.5), Inches(2), Inches(0.8))
        p.text_frame.text = txt
        p.fill.solid()
        p.fill.fore_color.rgb = RGBColor(30, 20, 50)
        p.line.color.rgb = ACCENT_PURPLE

    # ================= SLIDE 7: RIPES INTRO =================
    s7 = add_slide()
    add_title(s7, "Layer 1: Ripes (RISC-V Simulation)")
    
    # Left Box
    tf_l = add_text_box(s7, "What is Ripes?", Inches(1), Inches(2), Inches(5), Inches(4))
    tf_l.paragraphs[0].font.size = Pt(24)
    tf_l.paragraphs[0].font.color.rgb = ACCENT_RED
    
    bullets_l = [
        "Visual RISC-V Processor Simulator",
        "Pipeline visualization (IF, ID, EX, MEM, WB)",
        "Register file inspection",
        "Memory-Mapped I/O support"
    ]
    for b in bullets_l: add_bullet(tf_l, b)
    
    # Right Box
    tf_r = add_text_box(s7, "Why RISC-V?", Inches(7), Inches(2), Inches(5), Inches(4))
    tf_r.paragraphs[0].font.size = Pt(24)
    tf_r.paragraphs[0].font.color.rgb = ACCENT_RED
    
    bullets_r = [
        "Open-source ISA (no licensing)",
        "Reduced instruction set = efficiency",
        "Deterministic timing for safety",
        "Educational transparency"
    ]
    for b in bullets_r: add_bullet(tf_r, b)

    # ================= SLIDE 8: MEMORY MAPPED I/O =================
    s8 = add_slide()
    add_title(s8, "Memory-Mapped I/O Configuration")
    
    # Header
    add_text_box(s8, "Address              Name                 Bit Layout                         Purpose", Inches(1), Inches(2), Inches(11), Inches(0.5), font_size=16, color=ACCENT_BLUE, bold=True)
    
    # Rows
    rows = [
        "0xF0000000        Switch Bank 0      Bit 0: Smoke | 1-7: PM2.5     Read Sensor States",
        "0xF0000004        Switch Bank 1      Bits 0-7: CO2 Level              Read CO2 Sensor",
        "0xF0000008        LED Matrix         35x25 Pixel Grid (875 px)      Visual Feedback"
    ]
    
    y = 2.8
    for r in rows:
        add_text_box(s8, r, Inches(1), Inches(y), Inches(11), Inches(0.5), font_size=16, color=TEXT_MAIN)
        y += 0.8
        
    # Legend
    leg_y = 5.5
    add_text_box(s8, "LED Status Codes:", Inches(1), Inches(leg_y), Inches(4), Inches(0.5), font_size=18, color=TEXT_SEC)
    
    colors = [("GREEN = Safe", ACCENT_GREEN), ("RED = Danger", ACCENT_RED), ("BLUE = Fan Active", ACCENT_BLUE)]
    for i, (txt, col) in enumerate(colors):
        b = s8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1 + i*3.5), Inches(leg_y + 0.5), Inches(0.3), Inches(0.3))
        b.fill.solid()
        b.fill.fore_color.rgb = col
        add_text_box(s8, txt, Inches(1.4 + i*3.5), Inches(leg_y + 0.4), Inches(2), Inches(0.5), 16, TEXT_MAIN)

    # ================= SLIDE 9: ASSEMBLY CODE =================
    s9 = add_slide()
    add_title(s9, "RISC-V Assembly: Core Logic")
    
    code = """# Sensor Polling & Smoke Check (Critical Priority)
main_loop:
    li t0, 0xF0000000      # Load Switch Bank 0 Address
    lw s0, 0(t0)           # Read Raw Input (Smoke + PM2.5)
    
    li t1, 1               # Mask for Bit 0
    and a0, s0, t1         # Isolate Smoke Bit
    bnez a0, unsafe_mode   # BRANCH IMMEDIATE if Smoke!

# PM2.5 Parsing
    srli a1, s0, 1         # Shift right (remove Smoke bit)
    li t1, 0x7F            # Mask 0111 1111 (7-bit PM2.5)
    and a1, a1, t1         # Clean PM2.5 value
    
    li t2, 100             # Threshold constant
    bgt a1, t2, unsafe_mode"""
    
    add_code_block(s9, code, Inches(1), Inches(2), Inches(11), Inches(4.5))

    # ================= SLIDE 10: WOKWI INTRO =================
    s10 = add_slide()
    add_title(s10, "Layer 2: Wokwi (IoT Hardware Simulation)")
    
    # Hardware
    tf_h = add_text_box(s10, "Hardware Components", Inches(1), Inches(2), Inches(5), Inches(4))
    tf_h.paragraphs[0].font.size = Pt(22)
    tf_h.paragraphs[0].font.color.rgb = ACCENT_BLUE
    
    h_list = ["ESP32 Microcontroller (WiFi)", "Potentiometers (Analog Sensors)", "Slide Switch (Digital Smoke)", "RGB LEDs (Status)"]
    for i in h_list: add_bullet(tf_h, i)
    
    # Connectivity
    tf_c = add_text_box(s10, "Connectivity", Inches(7), Inches(2), Inches(5), Inches(4))
    tf_c.paragraphs[0].font.size = Pt(22)
    tf_c.paragraphs[0].font.color.rgb = ACCENT_BLUE
    
    c_list = ["WiFi: Wokwi-GUEST", "MQTT: broker.hivemq.com", "Topic: smart-corridor/sensors", "Topic: smart-corridor/commands"]
    for i in c_list: add_bullet(tf_c, i)

    # ================= SLIDE 11: WOKWI LOGIC =================
    s11 = add_slide()
    add_title(s11, "RISC-V Style Decision Logic (C++)")
    
    code_cpp = """// Priority-based hazard detection
if (effective_smoke > 30) {
    status_msg = "CRITICAL: SMOKE DETECTED";
    alert_level = "critical";      // Priority 1
}
else if (effective_temp > 35) {
    status_msg = "DANGER: HIGH TEMPERATURE";
    alert_level = "danger";        // Priority 2
}
else if (effective_pm25 > 80) {
    status_msg = "WARNING: HIGH PM2.5";
    alert_level = "warning";       // Priority 3
}"""
    add_code_block(s11, code_cpp, Inches(1), Inches(2), Inches(11), Inches(4.5))

    # ================= SLIDE 12: 3D SIM INTRO =================
    s12 = add_slide()
    add_title(s12, "Layer 3: 3D Digital Twin Simulation")
    
    # Tech Stack
    tf_t = add_text_box(s12, "Technology Stack", Inches(1), Inches(2), Inches(5), Inches(4))
    tf_t.paragraphs[0].font.size = Pt(22)
    tf_t.paragraphs[0].font.color.rgb = ACCENT_GREEN
    
    t_list = ["React (Component UI)", "Three.js / Fiber (3D Rendering)", "Rapier (Physics Engine)", "Vite (Build Tool)"]
    for i in t_list: add_bullet(tf_t, i)
    
    # Features
    tf_f = add_text_box(s12, "Environment Features", Inches(7), Inches(2), Inches(5), Inches(4))
    tf_f.paragraphs[0].font.size = Pt(22)
    tf_f.paragraphs[0].font.color.rgb = ACCENT_GREEN
    
    f_list = ["First-Person Controller", "Dynamic Smoke Particles", "Interactive Appliances", "Real-time Alert Lighting"]
    for i in f_list: add_bullet(tf_f, i)

    # ================= SLIDE 13: SIM FEATURES =================
    s13 = add_slide()
    add_title(s13, "Simulation Capabilities")
    
    caps = [
        ("🌫️ Fog System", "Thousands of particles visualizing smoke spread"),
        ("📊 Dashboard", "Live overlay showing CO2, PM2.5, AQI"),
        ("💨 Ventilation", "Fans spin and clear smoke on command"),
        ("🚨 Alert Lights", "House lights change color (Green -> Red)")
    ]
    
    for i, (t, d) in enumerate(caps):
        x = Inches(1 + (i%2)*6)
        y = Inches(2 + (i//2)*2.5)
        box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(20, 30, 25)
        box.line.color.rgb = ACCENT_GREEN
        
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = t
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = d
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_SEC

    # ================= SLIDE 14: MQTT =================
    s14 = add_slide()
    add_title(s14, "System Integration via MQTT")
    
    # Diagram
    # 3D Sim
    sim = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(3), Inches(2))
    sim.text_frame.text = "3D Simulation\n(React)\n[Publishes Sensors]"
    sim.fill.solid()
    sim.fill.fore_color.rgb = RGBColor(20, 40, 30)
    sim.line.color.rgb = ACCENT_GREEN
    
    # Broker
    broker = s14.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(5), Inches(2), Inches(3), Inches(4))
    broker.text_frame.text = "MQTT Broker\nhivemq.com\n\n<Topics>\n/sensors\n/commands"
    broker.fill.solid()
    broker.fill.fore_color.rgb = RGBColor(40, 30, 60)
    broker.line.color.rgb = ACCENT_PURPLE
    
    # Wokwi
    wok = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(3), Inches(3), Inches(2))
    wok.text_frame.text = "Wokwi (ESP32)\n(RISC-V Logic)\n[Publishes Commands]"
    wok.fill.solid()
    wok.fill.fore_color.rgb = RGBColor(20, 30, 50)
    wok.line.color.rgb = ACCENT_BLUE
    
    # Arrows
    s14.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4), Inches(3.5), Inches(5), Inches(3.5))
    s14.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8), Inches(3.5), Inches(9), Inches(3.5))

    # ================= SLIDE 15: IMPLEMENTATION FLOW =================
    s15 = add_slide()
    add_title(s15, "Implementation Flow")
    
    steps = [
        "Requirement Analysis: Define thresholds & safety standards",
        "RISC-V Logic Design: Assembly code in Ripes",
        "IoT Layer Development: ESP32 firmware in Wokwi",
        "3D Environment: React Three Fiber house model",
        "Integration: MQTT bridge connection",
        "Testing: Validation of all fail-safes"
    ]
    
    for i, step in enumerate(steps):
        box = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2 + i*0.8), Inches(9), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(25, 25, 35)
        box.line.color.rgb = ACCENT_BLUE
        
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {step}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MAIN

    # ================= SLIDE 16: TESTING =================
    s16 = add_slide()
    add_title(s16, "Testing & Validation Strategy")
    
    # Unit
    tf_u = add_text_box(s16, "Unit Testing", Inches(1), Inches(2), Inches(5), Inches(4))
    tf_u.paragraphs[0].font.size = Pt(22)
    tf_u.paragraphs[0].font.color.rgb = ACCENT_PURPLE
    
    u_list = ["Masking Logic (AND ops)", "Branch Logic (Thresholds)", "LED Drawing (Mem Addresses)"]
    for i in u_list: add_bullet(tf_u, i)
    
    # Integration
    tf_i = add_text_box(s16, "Integration Testing", Inches(7), Inches(2), Inches(5), Inches(4))
    tf_i.paragraphs[0].font.size = Pt(22)
    tf_i.paragraphs[0].font.color.rgb = ACCENT_PURPLE
    
    i_list = ["Smoke Switch Latency (<1 cycle)", "Noisy Data Rejection", "MQTT Round-Trip (<500ms)"]
    for i in i_list: add_bullet(tf_i, i)

    # ================= SLIDE 17: DEMO =================
    s17 = add_slide()
    add_title(s17, "Live Demo Phase")
    
    demo_box = s17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.5), Inches(9.33), Inches(3))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = RGBColor(30, 40, 30)
    demo_box.line.color.rgb = ACCENT_GREEN
    
    tf = demo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Demonstrating the full loop:"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\n1. Ripes: Real-time Register Switching\n2. Wokwi: Sensor & Actuator Hardware\n3. 3D Sim: Visual Environmental Response"
    p2.font.size = Pt(20)
    p2.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 18: CONCLUSION =================
    s18 = add_slide()
    add_title(s18, "Conclusion & Future Work")
    
    # Achievements
    tf_a = add_text_box(s18, "Key Achievements", Inches(1), Inches(2), Inches(5), Inches(3))
    tf_a.paragraphs[0].font.size = Pt(22)
    tf_a.paragraphs[0].font.color.rgb = ACCENT_GREEN
    
    a_list = ["Full RISC-V implementation", "IoT + Digital Twin integration", "<100ms Safety Response"]
    for i in a_list: add_bullet(tf_a, i)
    
    # Future
    tf_f = add_text_box(s18, "Future Enhancements", Inches(7), Inches(2), Inches(5), Inches(3))
    tf_f.paragraphs[0].font.size = Pt(22)
    tf_f.paragraphs[0].font.color.rgb = ACCENT_YELLOW
    
    f_list = ["FPGA Hardware Implementation", "Machine Learning Prediction", "Mobile App Monitoring"]
    for i in f_list: add_bullet(tf_f, i)
    
    # Thank You
    ty_box = s18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5.5), Inches(5.33), Inches(1.5))
    ty_box.fill.solid()
    ty_box.fill.fore_color.rgb = RGBColor(20, 20, 30)
    ty_box.line.color.rgb = ACCENT_BLUE
    
    tf_ty = ty_box.text_frame
    p = tf_ty.paragraphs[0]
    p.text = "Thank You!\nQuestions?"
    p.font.bold = True
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = TEXT_MAIN

    # SAVE
    output_file = "Smart_Corridor_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
